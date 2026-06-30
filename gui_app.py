import streamlit as st
from google import genai
from PIL import Image
import io

# Safe layout parsing library imports
try:
    from pdf2image import convert_from_bytes
except ImportError:
    st.error("Please ensure 'pdf2image' is included in your requirements.txt manifest.")

try:
    from pptx import Presentation
except ImportError:
    st.error("Please ensure 'python-pptx' is included in your requirements.txt manifest.")

# ==========================================================
# 🔑 API KEY CONFIGURATION
# ==========================================================
try:
    if "GEMINI_API_KEY" in st.secrets:
        MY_API_KEY = st.secrets["GEMINI_API_KEY"]
    else:
        MY_API_KEY = "YOUR_GEMINI_API_KEY_HERE"  # 
except Exception:
    MY_API_KEY = "YOUR_GEMINI_API_KEY_HERE"
# ==========================================================

# Initialize the Gemini Client
try:
    if MY_API_KEY == "YOUR_GEMINI_API_KEY_HERE" or not MY_API_KEY:
        st.error("⚠️ Missing API Key! Please configure your Gemini API key.")
        st.stop()
        
    client = genai.Client(api_key=MY_API_KEY)
except Exception as e:
    st.error(f"❌ Error initializing Gemini Client: {e}")
    st.stop()

# --- Page UI Configuration ---
st.set_page_config(page_title="Smart OCR Analyzer", page_icon="📊", layout="wide")

# Custom CSS for absolute responsiveness, big text, layout centering, and the isolated RED action button
st.markdown("""
    <style>
    /* Centering the main header and description */
    .centered-title {
        text-align: center;
        font-size: calc(2.2rem + 1.5vw) !important;
        font-weight: bold;
        margin-bottom: 0.5rem;
        color: inherit !important; 
    }
    .centered-subtitle {
        text-align: center;
        color: #A0A0A0;
        font-size: calc(1.0rem + 0.3vw) !important;
        margin-bottom: 2rem;
    }
    
    /* Enlarge the section labels */
    .big-label {
        font-size: 1.4rem !important;
        font-weight: bold !important;
        margin-bottom: 10px !important;
    }

    /* Restricting maximum image preview height so it doesn't break mobile workflows */
    .stImage img {
        max-height: 300px !important;
        width: auto !important;
        margin-left: auto;
        margin-right: auto;
        display: block;
        border-radius: 8px;
    }

    /* Style ONLY the standalone main Execute button to be red. 
       Excludes choice buttons sitting inside grid columns. */
    div.stButton > button:not(div[data-testid="stColumn"] button) {
        background-color: #D32F2F !important; /* Rich Crimson Red */
        color: white !important;
        font-weight: bold !important;
        font-size: 1.25rem !important;
        padding: 0.6rem 1rem !important;
        border: none !important;
        border-radius: 8px !important;
        transition: background-color 0.3s ease, transform 0.1s ease !important;
    }
    
    /* Hover state modification strictly for the standalone red execution button */
    div.stButton > button:not(div[data-testid="stColumn"] button):hover {
        background-color: #B71C1C !important; /* Deeper cherry red on hover */
        color: white !important;
        border: none !important;
    }
    
    /* Style optimization for text tracking within choice buttons */
    div[data-testid="stColumn"] div.stButton > button {
        white-space: pre-wrap !important;
    }

    /* Custom Footer Styling */
    .footer {
        text-align: center;
        padding: 20px;
        font-size: 0.9rem;
        color: #888888;
        margin-top: 50px;
    }
    .footer a {
        color: #D32F2F !important; /* Matches your red theme */
        text-decoration: none;
        font-weight: bold;
    }
    .footer a:hover {
        text-decoration: underline !important;
    }
    </style>
""", unsafe_allow_html=True)

# --- Centralized Title Banner ---
st.markdown('<div class="centered-title">📊 Smart OCR & File Analyzer</div>', unsafe_allow_html=True)
st.markdown('<div class="centered-subtitle">Upload any document, presentation, or image to cleanly extract structured context via Gemini 2.5 Flash</div>', unsafe_allow_html=True)
st.markdown("---")

# --- FIRST-TIME WELCOME BLOCK (Classic Expander Design) ---
if 'first_time_load' not in st.session_state:
    st.session_state.first_time_load = True

if st.session_state.first_time_load:
    st.toast("👋 Welcome to Smart OCR & File Analyzer!")
    
    with st.expander("🚀 Quick Intro: What does this application do?", expanded=True):
        st.markdown("""
        **Welcome!** This utility combines **Python & Streamlit** with the advanced vision capabilities of the **Gemini 2.5 Flash** model.
        
        * **🔍 Extract Text:** Instantly read text from images or PDFs (Optical Character Recognition).
        * **📝 Get Explanations:** Convert confusing forms or charts into plain, readable summaries.
        * **🔬 Deep Visual Audits:** Automatically flag structural components like signatures, graphics, or logos.
        
        *Simply upload your file on the left, choose your analysis options, and hit the big red button!*
        """)
    st.session_state.first_time_load = False

# --- Initialize Session States for Selection ---
if 'opt_ocr' not in st.session_state: st.session_state.opt_ocr = True
if 'opt_explain' not in st.session_state: st.session_state.opt_explain = False
if 'opt_deep' not in st.session_state: st.session_state.opt_deep = False

# --- UI Layout Grid Setup ---
col_input, col_output = st.columns([1, 1.2], gap="large")

with col_input:
    st.markdown('<p class="big-label">📁 Drag & Drop Document</p>', unsafe_allow_html=True)
    uploaded_file = st.file_uploader(
        "Upload Image, PDF, or PowerPoint", 
        type=['png', 'jpg', 'jpeg', 'webp', 'bmp'],
        label_visibility="collapsed"
    )
    
    st.markdown("---")
    st.markdown('<p class="big-label">⚙️ Analysis Options</p>', unsafe_allow_html=True)
    st.caption("Click the giant buttons below to select/deselect features before analyzing:")

    # High-fidelity, macro-sized selectable button rows
    btn_col1, btn_col2, btn_col3 = st.columns(3)
    
    with btn_col1:
        if st.button(
            f"🔍 OCR Text\n[{'ACTIVE' if st.session_state.opt_ocr else 'OFF'}]", 
            use_container_width=True, 
            type="primary" if st.session_state.opt_ocr else "secondary"
        ):
            st.session_state.opt_ocr = not st.session_state.opt_ocr
            st.rerun()

    with btn_col2:
        if st.button(
            f"📝 Explanation\n[{'ACTIVE' if st.session_state.opt_explain else 'OFF'}]", 
            use_container_width=True, 
            type="primary" if st.session_state.opt_explain else "secondary"
        ):
            st.session_state.opt_explain = not st.session_state.opt_explain
            st.rerun()

    with btn_col3:
        if st.button(
            f"🔬 Deep Analysis\n[{'ACTIVE' if st.session_state.opt_deep else 'OFF'}]", 
            use_container_width=True, 
            type="primary" if st.session_state.opt_deep else "secondary"
        ):
            st.session_state.opt_deep = not st.session_state.opt_deep
            st.rerun()

    # Extra spacing before execution
    st.markdown(" ")
    analyze_button = st.button("🚀 EXECUTE SMART ANALYSIS", use_container_width=True)

with col_output:
    st.markdown('<p class="big-label">🖥️ Output Preview Window</p>', unsafe_allow_html=True)
    
    if uploaded_file is not None:
        file_name = uploaded_file.name.lower()
        processing_data = None
        is_pptx = file_name.endswith('.pptx')
        is_pdf = file_name.endswith('.pdf')
        
        if not is_pdf and not is_pptx:
            # Native Image Stream
            img_preview = Image.open(uploaded_file)
            st.image(img_preview, caption=f"Active Target: {uploaded_file.name}")
            processing_data = img_preview
        elif is_pdf:
            st.info(f"📄 Active PDF Document: **{uploaded_file.name}**")
            try:
                # Convert first page of PDF to image format for spatial layout tracking
                pdf_images = convert_from_bytes(uploaded_file.read(), first_page=1, last_page=1)
                if pdf_images:
                    processing_data = pdf_images[0]
                    st.image(processing_data, caption="First Page Preview (Converted)", use_container_width=True)
            except Exception as pdf_err:
                st.warning("⚠️ Visual preview unavailable. Passing document payload directly.")
                uploaded_file.seek(0)
                processing_data = uploaded_file.getvalue()
        elif is_pptx:
            st.info(f"📊 Active PowerPoint Presentation: **{uploaded_file.name}**")
            try:
                prs = Presentation(uploaded_file)
                pptx_text_parts = []
                for i, slide in enumerate(prs.slides, 1):
                    pptx_text_parts.append(f"--- Slide {i} ---")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            pptx_text_parts.append(shape.text.strip())
                
                processing_data = "\n".join(pptx_text_parts)
                st.success(f"✅ Successfully parsed {len(prs.slides)} slides! Ready for processing.")
            except Exception as pptx_err:
                st.error(f"Failed to parse PowerPoint presentation: {pptx_err}")

        # Handle Action Execution
        if analyze_button and processing_data is not None:
            with st.spinner("🤖 Synthesizing file layout and text context..."):
                
                # Dynamically construct the systemic task prompt
                prompt_parts = []
                if st.session_state.opt_ocr:
                    prompt_parts.append("Perform OCR / Content Extraction: Read and extract all contents exactly as they appear. Maintain structure.")
                if st.session_state.opt_explain:
                    prompt_parts.append("Provide a clear, cohesive narrative explanation detailing the primary subject matter of this file.")
                if st.session_state.opt_deep:
                    prompt_parts.append("Identify layout and presentation aesthetics: Analyze structural topics, core insights, and provide a granular breakdown.")
                
                if not prompt_parts:
                    st.warning("⚠️ Please select at least one of the large analysis options block modules above!")
                    st.stop()
                    
                final_prompt = "\n\n".join(prompt_parts)

                try:
                    # Case 1: Structured string data (PowerPoint text maps)
                    if isinstance(processing_data, str):
                        full_content = f"{final_prompt}\n\nHere is the presentation content data text extracted directly from the slides:\n{processing_data}"
                        response = client.models.generate_content(
                            model='gemini-2.5-flash',
                            contents=full_content
                        )
                    # Case 2: Raw data bytes handling (PDF fallback) - Packaged securely for GenAI SDK
                    elif isinstance(processing_data, bytes):
                        response = client.models.generate_content(
                            model='gemini-2.5-flash',
                            contents=[
                                final_prompt,
                                {
                                    "mime_type": "application/pdf",
                                    "data": processing_data
                                }
                            ]
                        )
                    # Case 3: Images (Native image uploads or successful PDF conversion slides)
                    else:
                        response = client.models.generate_content(
                            model='gemini-2.5-flash',
                            contents=[processing_data, final_prompt]
                        )
                    
                    st.success("✅ File Successfully Analyzed!")
                    st.markdown("### 📝 Analysis Summary Breakdown:")
                    st.info(response.text)
                    
                except Exception as ai_err:
                    st.error(f"AI Generation Failed: {ai_err}")
    else:
        st.info("Awaiting file upload. Drop an item on the left column block to initiate data rendering.")

# --- Footer Section ---
st.markdown("---") 
footer_html = """
<div class="footer">
    © 2026 Afsal Hamad. All Rights Reserved. | 
    🌐 Source Code: <a href="https://github.com/afsalahamad/Python-based-Analyzer" target="_blank">GitHub Repository</a>
</div>
"""
st.markdown(footer_html, unsafe_allow_html=True)